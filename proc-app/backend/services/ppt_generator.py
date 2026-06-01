"""PowerPoint deck generator.

Generates a findings + KPI deck from a KPI-dashboard run output.
Uses python-pptx with a built-in clean layout — no external template.
"""
from __future__ import annotations

import io
from datetime import datetime
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


# Procvault indigo palette
INDIGO_900 = RGBColor(0x1E, 0x1B, 0x4B)
INDIGO_700 = RGBColor(0x43, 0x38, 0xCA)
INDIGO_500 = RGBColor(0x6E, 0x59, 0xF7)
INDIGO_50 = RGBColor(0xEE, 0xEC, 0xFE)
INK_900 = RGBColor(0x0F, 0x14, 0x2B)
INK_600 = RGBColor(0x4B, 0x55, 0x63)
INK_400 = RGBColor(0x96, 0x9F, 0xB0)
SUCCESS = RGBColor(0x15, 0x80, 0x4A)
WARN = RGBColor(0xC0, 0x73, 0x12)
DANGER = RGBColor(0xB3, 0x26, 0x2D)


PILLAR_LABELS = {
    "op-model": "Operating Model",
    "buying-channel": "Buying Channel",
    "org-structure": "Organisation Structure",
    "doa": "Delegation of Authority",
}

STATUS_COLOR = {
    "in": SUCCESS, "under": WARN, "over": DANGER, "unknown": INK_400,
}


# ============================================================================
# Public entry points
# ============================================================================

def generate_findings_deck(engagement: dict, kpi_dashboard: dict) -> bytes:
    """Build a multi-slide deck covering exec summary + per-pillar findings."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _slide_cover(prs, engagement)
    _slide_exec_summary(prs, engagement, kpi_dashboard)
    _slide_pillar_overview(prs, kpi_dashboard)

    for pid in ["op-model", "buying-channel", "org-structure", "doa"]:
        _slide_pillar_detail(prs, pid, kpi_dashboard)

    _slide_top_alerts(prs, kpi_dashboard)
    _slide_appendix(prs, engagement, kpi_dashboard)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_exec_summary_deck(engagement: dict, kpi_dashboard: dict) -> bytes:
    """Single-slide exec summary suitable for executive briefing."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _slide_cover(prs, engagement)
    _slide_exec_summary(prs, engagement, kpi_dashboard)
    _slide_pillar_overview(prs, kpi_dashboard)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ============================================================================
# Slide builders
# ============================================================================

def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_text(slide, left, top, width, height, text, *, size=14, bold=False,
              color=INK_900, italic=False, font="Geist"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return tb


def _add_rect(slide, left, top, width, height, fill=INDIGO_500, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    return sh


def _slide_cover(prs, engagement):
    slide = _blank(prs)
    _add_rect(slide, 0, 0, 13.33, 7.5, fill=INDIGO_900)
    _add_rect(slide, 0, 6.5, 13.33, 1.0, fill=INDIGO_700)

    _add_text(slide, 0.6, 1.6, 12, 0.5, "PROCVAULT", size=12, color=INDIGO_50)
    _add_text(slide, 0.6, 2.0, 12, 1.5,
              "Procurement Maturity Assessment",
              size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    _add_text(slide, 0.6, 3.5, 12, 1.0,
              engagement.get("client_name", "—"),
              size=28, italic=True, color=INDIGO_50, font="Newsreader")
    sub = f"{engagement.get('industry','').title()} · " \
          f"{(engagement.get('sub_segment') or '').replace('_',' ').title()}"
    _add_text(slide, 0.6, 4.5, 12, 0.5, sub, size=16, color=INDIGO_50)

    plants = engagement.get("plants") or []
    facts = f"{len(plants)} plants · ₹{engagement.get('annual_spend_inr_cr','—')} Cr spend · " \
            f"{engagement.get('fte_count','—')} FTEs"
    _add_text(slide, 0.6, 5.2, 12, 0.4, facts, size=13, color=INDIGO_50)

    _add_text(slide, 0.6, 6.7, 12, 0.5,
              f"Generated {datetime.now().strftime('%d %B %Y')}",
              size=10, color=INDIGO_50)


def _slide_exec_summary(prs, engagement, kp):
    slide = _blank(prs)
    _add_text(slide, 0.5, 0.4, 12, 0.5, "EXECUTIVE SUMMARY", size=11, color=INK_600)
    _add_text(slide, 0.5, 0.85, 12, 0.8, "How the function performs against benchmarks",
              size=26, bold=True)

    pillars = list(kp.get("pillar_summary", {}).items())
    # Skip pillars where score is None (needs_qre stub, or otherwise unscored)
    scored = [p for p in pillars if (p[1].get("pillar_score") or {}).get("score") is not None]
    avg = (sum(p[1]["pillar_score"]["score"] for p in scored) / len(scored)) if scored else 0
    label = _maturity_label(avg)

    # Headline
    headline = (f"{engagement.get('client_name','—')} sits at {label} "
                f"({avg:.1f}/5) across {len(pillars)} pillars.")
    _add_text(slide, 0.5, 2.0, 12, 1.0, headline, size=20, font="Newsreader", italic=True)

    in_band  = sum(1 for k in kp.get("kpis", []) if k.get("status") == "in")
    under    = sum(1 for k in kp.get("kpis", []) if k.get("status") == "under")
    over     = sum(1 for k in kp.get("kpis", []) if k.get("status") == "over")
    total    = len(kp.get("kpis", []))

    para = (f"Of {total} benchmarked KPIs, {in_band} sit inside their band, "
            f"{under} below, and {over} above. {under + over} require attention.")
    _add_text(slide, 0.5, 3.2, 12, 1.0, para, size=15, color=INK_600)

    # KPI counter strip
    metrics = [
        ("In band", in_band, SUCCESS),
        ("Below", under, WARN),
        ("Above", over, DANGER),
        ("Avg maturity", f"{avg:.1f}/5", INDIGO_700),
    ]
    for i, (lbl, v, col) in enumerate(metrics):
        x = 0.5 + i * 3.0
        _add_rect(slide, x, 4.6, 2.8, 1.8, fill=INDIGO_50)
        _add_text(slide, x + 0.2, 4.7, 2.6, 0.4, lbl, size=11, color=INK_600)
        _add_text(slide, x + 0.2, 5.1, 2.6, 1.0, str(v), size=32, bold=True, color=col)

    _add_text(slide, 0.5, 6.8, 12, 0.3,
              f"Total spend: ₹{kp.get('portfolio',{}).get('total_spend_inr_cr','—')} Cr · "
              f"{kp.get('portfolio',{}).get('po_count','—')} POs · "
              f"{kp.get('portfolio',{}).get('mg_count','—')} MGs",
              size=10, color=INK_400)


def _slide_pillar_overview(prs, kp):
    slide = _blank(prs)
    _add_text(slide, 0.5, 0.4, 12, 0.5, "PILLAR MATURITY", size=11, color=INK_600)
    _add_text(slide, 0.5, 0.85, 12, 0.8, "Per-pillar scores · weighted across themes",
              size=26, bold=True)

    pillars = list(kp.get("pillar_summary", {}).items())
    if not pillars:
        return
    col_w = 11.0 / len(pillars)
    for i, (pid, summary) in enumerate(pillars):
        x = 0.7 + i * col_w
        s = summary.get("pillar_score", {})
        score = s.get("score", 0)
        lbl = s.get("label", "")
        _add_rect(slide, x, 2.0, col_w - 0.3, 3.4, fill=INDIGO_50)
        _add_text(slide, x + 0.2, 2.1, col_w - 0.5, 0.4,
                  PILLAR_LABELS.get(pid, pid), size=12, bold=True, color=INK_600)
        _add_text(slide, x + 0.2, 2.6, col_w - 0.5, 1.4,
                  f"{score:.1f}" if isinstance(score, (int, float)) else str(score),
                  size=56, bold=True, color=INDIGO_700)
        _add_text(slide, x + 0.2, 4.1, col_w - 0.5, 0.4, lbl, size=12, color=INK_600)
        kc = summary.get("kpi_count", 0)
        ib = summary.get("in_band", 0)
        un = summary.get("under", 0)
        ov = summary.get("over", 0)
        _add_text(slide, x + 0.2, 4.6, col_w - 0.5, 0.4,
                  f"{kc} KPIs", size=11, bold=True, color=INK_600)
        _add_text(slide, x + 0.2, 4.95, col_w - 0.5, 0.4,
                  f"{ib} in · {un} below · {ov} above", size=10, color=INK_400)


def _slide_pillar_detail(prs, pid, kp):
    pillar_kpis = [k for k in kp.get("kpis", []) if k.get("pillar") == pid]
    if not pillar_kpis:
        return
    slide = _blank(prs)
    label = PILLAR_LABELS.get(pid, pid)
    summary = kp.get("pillar_summary", {}).get(pid, {})
    score = summary.get("pillar_score", {})

    _add_text(slide, 0.5, 0.4, 12, 0.5, label.upper(), size=11, color=INK_600)
    _add_text(slide, 0.5, 0.85, 12, 0.8,
              f"{label} — maturity {score.get('score', '—')} ({score.get('label', '')})",
              size=24, bold=True)

    # Top 6 KPIs sorted: outliers first
    sev_rank = {"over": 0, "under": 1, "unknown": 2, "in": 3}
    pillar_kpis.sort(key=lambda k: sev_rank.get(k.get("status"), 9))
    rows = pillar_kpis[:6]

    top = 1.9
    h = 0.85
    for i, k in enumerate(rows):
        y = top + i * h
        col = STATUS_COLOR.get(k.get("status"), INK_400)
        _add_rect(slide, 0.5, y, 0.08, h - 0.1, fill=col)
        _add_text(slide, 0.7, y, 6.5, 0.4, k.get("label", ""), size=13, bold=True)
        finding = (k.get("finding") or "")[:140]
        _add_text(slide, 0.7, y + 0.4, 6.5, 0.4, finding, size=10, color=INK_600)
        value = k.get("value")
        unit = k.get("unit", "")
        _add_text(slide, 7.5, y + 0.05, 2.2, 0.5,
                  f"{value:.1f}" if isinstance(value, (int, float)) else str(value),
                  size=20, bold=True)
        _add_text(slide, 7.5, y + 0.55, 2.2, 0.3, unit, size=10, color=INK_400)
        band = k.get("band", {})
        delta = k.get("delta", "")
        _add_text(slide, 9.8, y + 0.05, 3.3, 0.4,
                  f"Band: {band.get('low','—')}–{band.get('high','—')}", size=11, color=INK_600)
        _add_text(slide, 9.8, y + 0.45, 3.3, 0.4, delta, size=10, color=col)


def _slide_top_alerts(prs, kp):
    slide = _blank(prs)
    _add_text(slide, 0.5, 0.4, 12, 0.5, "ATTENTION", size=11, color=INK_600)
    _add_text(slide, 0.5, 0.85, 12, 0.8, "Top KPIs outside benchmark band",
              size=26, bold=True)

    outliers = [k for k in kp.get("kpis", []) if k.get("status") in ("under", "over")]
    sev_rank = {"over": 0, "under": 1}
    outliers.sort(key=lambda k: sev_rank.get(k.get("status"), 9))
    top = outliers[:7]
    if not top:
        _add_text(slide, 0.5, 2.0, 12, 1, "All KPIs sit inside their benchmark band.",
                  size=18, italic=True, font="Newsreader", color=SUCCESS)
        return

    y = 1.8
    for k in top:
        col = STATUS_COLOR.get(k.get("status"), INK_400)
        _add_rect(slide, 0.5, y, 0.08, 0.6, fill=col)
        _add_text(slide, 0.7, y, 8.5, 0.4,
                  f"{PILLAR_LABELS.get(k['pillar'], k['pillar'])} · {k.get('label', '')}",
                  size=12, bold=True)
        _add_text(slide, 0.7, y + 0.35, 8.5, 0.3,
                  (k.get("finding") or "")[:160], size=10, color=INK_600)
        _add_text(slide, 9.4, y + 0.05, 1.6, 0.4,
                  f"{k.get('value'):.1f}" if isinstance(k.get('value'), (int, float)) else str(k.get('value', '—')),
                  size=16, bold=True, color=col)
        _add_text(slide, 9.4, y + 0.45, 1.6, 0.3, k.get("unit", ""),
                  size=9, color=INK_400)
        _add_text(slide, 11.0, y + 0.05, 2.1, 0.4,
                  k.get("delta", ""), size=10, color=col)
        y += 0.75


def _slide_appendix(prs, engagement, kp):
    slide = _blank(prs)
    _add_text(slide, 0.5, 0.4, 12, 0.5, "APPENDIX", size=11, color=INK_600)
    _add_text(slide, 0.5, 0.85, 12, 0.8, "Methodology + citations",
              size=26, bold=True)
    body = (
        "• Pillar scores cascade through Function defaults → Industry overlays → Engagement overrides.\n"
        "• Every KPI carries source, year, and confidence (see drill-down in the live dashboard).\n"
        f"• Data window: top {kp.get('mg_count', '—')} material groups by spend.\n"
        f"• Engine version: Procvault V1.\n"
        f"• Generated: {datetime.now().strftime('%d %B %Y, %H:%M IST')}."
    )
    _add_text(slide, 0.5, 2.0, 12, 4.0, body, size=14, color=INK_600)


def _maturity_label(score):
    if score < 1.5: return "Initial"
    if score < 2.5: return "Developing"
    if score < 3.5: return "Defined"
    if score < 4.5: return "Managed"
    return "Optimised"
